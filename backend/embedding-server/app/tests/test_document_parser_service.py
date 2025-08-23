import base64
import io
import os
from email.message import EmailMessage

import pytest

from app.services.document_parser_service import DocumentParserService


@pytest.fixture()
def parser():
    # Disable OCR for deterministic tests
    return DocumentParserService(ocr_enabled=False)


def _write(tmp_path, name: str, data: bytes) -> str:
    p = tmp_path / name
    p.write_bytes(data)
    return str(p)


def test_parse_txt(parser, tmp_path):
    path = _write(tmp_path, "sample.txt", b"Hello TXT")
    docs = parser.parse_file(path)
    assert any("Hello TXT" in d["content"] for d in docs)


def test_parse_md(parser, tmp_path):
    path = _write(tmp_path, "sample.md", b"# Title\n\nHello MD")
    docs = parser.parse_file(path)
    assert any("Hello MD" in d["content"] for d in docs)


def test_parse_html(parser, tmp_path):
    html = b"<html><body><p>Hello HTML</p></body></html>"
    path = _write(tmp_path, "sample.html", html)
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Hello HTML" in joined


def test_parse_json(parser, tmp_path):
    path = _write(tmp_path, "data.json", b"{\n  \"a\": 1\n}")
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "a" in joined


def test_parse_yaml(parser, tmp_path):
    path = _write(tmp_path, "data.yaml", b"a: 1\n")
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "a" in joined


def test_parse_xml(parser, tmp_path):
    path = _write(tmp_path, "data.xml", b"<root><p>Hello XML</p></root>")
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Hello XML" in joined


def test_parse_csv(parser, tmp_path):
    path = _write(tmp_path, "data.csv", b"name,age\nAlice,30\nBob,40\n")
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "name" in joined and "Alice" in joined


def test_parse_eml(parser, tmp_path):
    msg = EmailMessage()
    msg["Subject"] = "Test Email"
    msg["From"] = "a@example.com"
    msg["To"] = "b@example.com"
    msg.set_content("Hello EML")
    path = _write(tmp_path, "mail.eml", msg.as_bytes())
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Test Email" in joined and "Hello EML" in joined


def test_parse_docx(parser, tmp_path):
    try:
        import docx  # type: ignore
    except Exception:
        pytest.skip("python-docx not available")
    doc = docx.Document()
    doc.add_paragraph("Hello DOCX")
    bio = io.BytesIO()
    doc.save(bio)
    path = _write(tmp_path, "sample.docx", bio.getvalue())
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Hello DOCX" in joined


def test_parse_pptx(parser, tmp_path):
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        pytest.skip("python-pptx not available")
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Title Only / Blank variant depending on template
    slide = prs.slides.add_slide(slide_layout)
    left = top = width = height = None
    # Add a textbox if shape factory exists; fallback: use title
    if hasattr(slide.shapes, "add_textbox"):
        txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        tf = txBox.text_frame
        tf.text = "Hello PPTX"
    else:
        slide.shapes.title.text = "Hello PPTX"  # type: ignore
    bio = io.BytesIO()
    prs.save(bio)
    path = _write(tmp_path, "sample.pptx", bio.getvalue())
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Hello PPTX" in joined


def test_parse_xlsx(parser, tmp_path):
    try:
        import pandas as pd  # type: ignore
    except Exception:
        pytest.skip("pandas not available")
    df = __import__("pandas").DataFrame({"name": ["Alice", "Bob"], "age": [30, 40]})  # type: ignore
    bio = io.BytesIO()
    try:
        with pd.ExcelWriter(bio, engine="openpyxl") as writer:  # type: ignore
            df.to_excel(writer, index=False, sheet_name="Sheet1")
    except Exception:
        pytest.skip("openpyxl not available or failed to write XLSX")
    path = _write(tmp_path, "sample.xlsx", bio.getvalue())
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Alice" in joined and "Bob" in joined


def test_parse_epub(parser, tmp_path):
    try:
        from ebooklib import epub  # type: ignore
    except Exception:
        pytest.skip("EbookLib not available")
    book = epub.EpubBook()
    book.set_identifier("id123")
    book.set_title("Sample Book")
    book.set_language("en")
    c1 = epub.EpubHtml(title="Intro", file_name="chap_01.xhtml", content="<h1>Hello EPUB</h1>")
    book.add_item(c1)
    book.toc = (epub.Link("chap_01.xhtml", "Intro", "intro"),)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", c1]
    bio = io.BytesIO()
    epub.write_epub(bio, book)  # type: ignore
    path = _write(tmp_path, "sample.epub", bio.getvalue())
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Hello EPUB" in joined


def test_parse_pdf(parser, tmp_path):
    # Generate a simple PDF if reportlab is available; otherwise skip
    try:
        from reportlab.pdfgen import canvas  # type: ignore
        from reportlab.lib.pagesizes import letter  # type: ignore
    except Exception:
        pytest.skip("reportlab not available to generate a test PDF")
    bio = io.BytesIO()
    c = canvas.Canvas(bio, pagesize=letter)
    c.drawString(100, 700, "Hello PDF")
    c.save()
    path = _write(tmp_path, "sample.pdf", bio.getvalue())
    docs = parser.parse_file(path)
    joined = "\n".join(d["content"] for d in docs)
    assert "Hello PDF" in joined


